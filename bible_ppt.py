import pickle
import re
from bible_function import (
    abbreviation,
    parse_paragraph,
    dict_contents,
    parsing_contents,
)
from pptx import Presentation
from pptx.util import Pt
from win32com.client import Dispatch
from os.path import abspath, join, dirname

BASE_DIR = dirname(abspath(__file__))


TEMPLATE_PATH = join(BASE_DIR, "verse_template.pptx")
SAVEFILE_NAME = join(BASE_DIR, "new-file-name.pptx")
# BIBLE_DF = join(BASE_DIR, "bible_dataframe.pickle")
# BIBLE_DICT = join(BASE_DIR, "bible_dict.pickle")

## Load pickle
with open("bible_list.pickle", "rb") as fr:
    bible_list = pickle.load(fr)
with open("bible_dict.pickle", "rb") as fr:
    bookDict = pickle.load(fr)

# PPTX 생성 함수 : 성경 본문(제목포함) + 한글 텍스트 파일
def present(
    resultList,
    template_path,
    save_file_name,
    main_book,
    main_chapter,
    main_verse_start,
    main_verse_end,
    keys,
    contentsDict,
    main_title,
):
    prs = Presentation(template_path)

    # 타이틀 삽입
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "성경봉독"
    subtitle.text = (
        f"{bookDict[main_book]} {main_chapter}장 {main_verse_start}-{main_verse_end}절"
    )
    # slide.shapes[2].text_frame.text = "(신약 p.)"

    # 본문 삽입
    for key in keys:
        contents_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(contents_slide_layout)
        title = slide.shapes.title
        contents = slide.placeholders[1]

        title.text = key
        contents.text = contentsDict[key]

    #############
    # 타이틀 삽입
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "말씀선포"
    subtitle.text = main_title

    for book, chapter, verse_start, verse_end, contents in resultList:
        if not verse_end:
            key = f"{bookDict[book]} {chapter}장 {verse_start}절"
            # pptx
            contents_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(contents_slide_layout)
            title_placeholder = slide.shapes.title
            contents_placeholder = slide.placeholders[1]

            title_placeholder.text = key
            contents_placeholder.text = contents[0]
        else:
            key = f"{bookDict[book]} {chapter}장 {verse_start}-{verse_end}절"

            contents_slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(contents_slide_layout)

            title_placeholder = slide.shapes.title
            title_placeholder.text = key

            contents_text_frame = slide.placeholders[1].text_frame
            for _ in range(len(contents) - 1):
                contents_text_frame.add_paragraph()
            for i in range(len(contents)):
                tmp1 = contents_text_frame.paragraphs[i].add_run()
                try:  # 본문에 실수로 구절번호를 생략한 경우 예외처리
                    tmp1.text = re.search("\d+", contents[i]).group()
                except:
                    pass
                tmp1.font.size = Pt(36)
                tmp2 = contents_text_frame.paragraphs[i].add_run()
                tmp2.text = re.sub("\d+", "", contents[i])

    prs.save(save_file_name)


# 파워포인트 열기
def open_ppt(filename=SAVEFILE_NAME):
    PPT = Dispatch("PowerPoint.Application")
    PPT.Visible = True
    PPT.Presentations.Open(abspath(filename))


# 성경 본문만 PPT 만들기
def make_ppt_only_text(paragraph):
    main_book, main_chapter, main_verse_start, main_verse_end = parse_paragraph(
        paragraph
    )
    keys, contentsDict = dict_contents(
        main_book, main_chapter, main_verse_start, main_verse_end
    )
    prs = Presentation(TEMPLATE_PATH)

    for key in keys:
        contents_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(contents_slide_layout)
        title = slide.shapes.title
        contents = slide.placeholders[1]

        title.text = key
        contents.text = contentsDict[key]

    prs.save(SAVEFILE_NAME)

    open_ppt(SAVEFILE_NAME)


def make_ppt_text(raw):
    (
        resultList,
        main_book,
        main_chapter,
        main_verse_start,
        main_verse_end,
        main_title,
    ) = parsing_contents(raw)

    print(main_book, main_chapter, main_verse_start, main_verse_end, main_title)

    keys, contentsDict = dict_contents(
        main_book, main_chapter, main_verse_start, main_verse_end
    )

    # print(resultList)

    present(
        resultList=resultList,
        template_path=TEMPLATE_PATH,
        save_file_name=SAVEFILE_NAME,
        main_book=main_book,
        main_chapter=main_chapter,
        main_verse_start=main_verse_start,
        main_verse_end=main_verse_end,
        keys=keys,
        contentsDict=contentsDict,
        main_title=main_title,
    )

    # open_ppt(SAVEFILE_NAME)


if __name__ == "__main__":
    make_ppt_text(raw)