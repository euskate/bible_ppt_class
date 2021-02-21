# 이전 목표 완성 : 찬양PPT 복사 배경이미지 소스 원본 붙여넣기
# 다음 목표
# 성경봉독 및 말씀선포 복사 붙여넣기 오류 잡기 (찬양 제목 부분)
# 찬양 슬라이드 삽입 위치 정상화

import sys
import win32com.client
from time import sleep
from os import listdir
import os
from bible_function import (
    abbreviation,
    parse_paragraph,
    dict_contents,
    parsing_contents,
    extract_main_verse,
    copy_contents,
    bookDict,
)
from bible_ppt import make_ppt_text, SAVEFILE_NAME

path = (
    "c:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\오전예배 (16x9)_2021_base.pptx"
)

hymn1_path = (
    "C:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\새찬송가16x9\\NHymn016h_Wide.PPT"
)


class Powerpoint:
    """
    Python PowerPoint win32 object
    """

    # 파워포인트 생성
    def init_app(self):
        self.app = win32com.client.Dispatch("PowerPoint.Application")
        self.app.Visible = True
        self.count = 0  # 전체 슬라이드 수 초기화
        return self.app

    # 프레젠테이션 열기
    def open_prs(self, path):
        self.path = path
        self.prs = self.app.Presentations.Open(self.path)
        return self.prs

    # 프레젠테이션 닫기
    def close_prs(self):
        self.prs.Close()

    # 전체 슬라이드 수 세기
    def get_count_slide(self):
        self.count = len(self.prs.Slides)  # 전체 슬라이드 길이를 반환한다

    # 전체 슬라이드 복사
    def copy_slide_all(self):
        self.get_count_slide()
        self.prs.Slides.Range(range(1, self.count + 1)).copy()

    # 슬라이드 복사
    def copy_slide(self, start_slide_number, end_slide_number):
        self.prs.Slides.Range(range(start_slide_number, end_slide_number + 1)).copy()

    # 슬라이드 붙여넣기
    def paste_slide(self, paste_slide_number):
        self.prs.Slides.paste(paste_slide_number)

    # 구역 가져오기
    def get_section_number(self, section_index):
        """[구역 가져오기]
        Args : section_index(int) : 몇번째 구역인지 설정
        Returns:
            [tuple]: (first_slide, section_count)
        """
        sec = self.prs.SectionProperties
        self.first_slide = sec.FirstSlide(section_index)  # 구역 첫번째 페이지
        self.section_count = sec.SlidesCount(section_index)  # 구역 슬라이드 개수
        return self.first_slide, self.section_count

    # 슬라이드 디자인 복붙
    def copy_desgin_slide(self, dst_prs, first_slide_number):
        """[슬라이드 디자인 복붙]

        Args:
            dst_prs ([Presentation_Object]): 붙여넣을 프레젠테이션
            first_slide_number ([int]): 붙여넣기 시작할 슬라이드 번호, self.get_section_number()에서 가져온다.
        """
        self.get_count_slide()  # 복사할 페이지 수를 가져온다.
        for i in range(self.count):
            dst_prs.Slides(first_slide_number + i).Design = self.prs.Slides(
                i + 1
            ).Design

    # 구역 지우기
    def del_section(self, first_slide, section_count):
        self.prs.Slides.Range(range(first_slide, first_slide + section_count)).Delete()

    # 찬양 번호 바꾸기
    def change_hymn_number(self, hymn_number, section_index):
        first_slide, _ = self.get_section_number(section_index)
        self.prs.Slides(first_slide).Shapes(2).Textframe.TextRange.Runs(
            2
        ).Text = hymn_number

    # 전환 애니메이션 적용하기
    def change_transition(self, start=1, end=1):
        """
        docstring
        """
        # 화면전환 : 밝기변화
        self.get_count_slide()
        end = self.count

        self.prs.Slides.Range(
            range(start, end)
        ).SlideShowTransition.EntryEffect = 3849  # 밝기변화
        self.prs.Slides.Range(
            range(start, end)
        ).SlideShowTransition.Duration = 0.5  # 시간
        pass

    def get_section(self, section_name):
        section_index = [
            i
            for i in range(1, self.prs.SectionProperties.Count + 1)
            if self.prs.SectionProperties.Name(i) == section_name
        ][0]
        sec = self.prs.SectionProperties
        first_slide = sec.FirstSlide(section_index)  # 구역 첫번째 페이지
        section_count = sec.SlidesCount(section_index)  # 구역 슬라이드 개수
        return first_slide, section_count

    # 제목 레이아웃 부제목 변경
    def change_subtitle(self, slide, subtitle):
        self.prs.Slides(slide).Shapes(2).Textframe.TextRange.Text = subtitle

    # '성경봉독' 섹션 입력
    def input_verse(self, slide, keys, contentsDict):
        for key in keys:
            slide += 1
            self.prs.Slides.AddSlide(
                slide, self.prs.SlideMaster.CustomLayouts(2)
            )  # 슬라이드 추가
            self.prs.Slides(
                slide
            ).Shapes.Title.TextFrame.TextRange.Text = key  # 제목 텍스트 입력
            self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Text = contentsDict[
                key
            ]  # 내용 텍스트 입력
        pass

    # '말씀 선포' 섹션 입력
    def input_hwp(self, slide, resultList):
        for book, chapter, verse_start, verse_end, contents in resultList:
            slide += 1  # 섹션 시작 슬라이드 다음부터
            if not verse_end:
                key = f"{bookDict[book]} {chapter}장 {verse_start}절"
                # pptx
                self.prs.Slides.AddSlide(slide, self.prs.SlideMaster.CustomLayouts(2))
                self.prs.Slides(
                    slide
                ).Shapes.Title.TextFrame.TextRange.Text = key  # 제목 텍스트
                self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Text = contents[
                    0
                ]  # 내용 텍스트

            else:
                key = f"{bookDict[book]} {chapter}장 {verse_start}-{verse_end}절"

                self.prs.Slides.AddSlide(slide, self.prs.SlideMaster.CustomLayouts(3))
                self.prs.Slides(
                    slide
                ).Shapes.Title.TextFrame.TextRange.Text = key  # 제목 텍스트
                self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Text = "\n".join(
                    contents
                )  # 내용 텍스트
                self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Runs(
                    1
                ).Font.Size = 28  # 구절 폰트 크기 변경
                start_chr = 0
                for _ in range(len(contents) - 1):
                    start_chr = (
                        self.prs.Slides(slide)
                        .Shapes(2)
                        .TextFrame.TextRange.Find("\n", start_chr)
                        .Start
                    )  # 문단 넘김 시작점 찾기
                    self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Characters(
                        start_chr + 1, 9999
                    ).Runs(
                        1
                    ).Font.Size = 28  # 구절 폰트 크기 변경

    """
    # '말씀 선포' 섹션 입력
    def input_hwp(self, slide, resultList):
        for i, (subtitle, text) in enumerate(resultList):
            slide += 1  # 섹션 시작 슬라이드 다음부터
            if len(text) < 2:
                self.prs.Slides.AddSlide(slide, self.prs.SlideMaster.CustomLayouts(2))
                self.prs.Slides(
                    slide
                ).Shapes.Title.TextFrame.TextRange.Text = resultList[i][
                    0
                ]  # 제목 텍스트 입력
                self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Text = resultList[
                    i
                ][1][
                    0
                ]  # 내용 텍스트 입력
            else:
                self.prs.Slides.AddSlide(slide, self.prs.SlideMaster.CustomLayouts(3))
                self.prs.Slides(
                    slide
                ).Shapes.Title.TextFrame.TextRange.Text = resultList[i][
                    0
                ]  # 제목 텍스트
                self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Text = "\n".join(
                    resultList[i][1]
                )  # 내용 텍스트
                self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Runs(
                    1
                ).Font.Size = 36  # 구절 폰트 크기 변경
                for _ in range(len(text) - 1):
                    start_chr = (
                        self.prs.Slides(slide)
                        .Shapes(2)
                        .TextFrame.TextRange.Find("\n")
                        .Start
                    )  # 문단 넘김 시작점 찾기
                    self.prs.Slides(slide).Shapes(2).TextFrame.TextRange.Characters(
                        start_chr + 1, 9999
                    ).Runs(
                        1
                    ).Font.Size = 36  # 구절 폰트 크기 변경 */
    """


### 한글 ###
hwp_path = "D:\\Documents\\카카오톡 받은 파일\\01. 예수를 가르쳐 복음을 전하니(ppt용 성경구절).hwp"


def hwp(hwp_path):
    """ 한글함수 """
    import win32clipboard

    hwp = win32com.client.gencache.EnsureDispatch("HWPFrame.HwpObject")
    hwp.RegisterModule("FilePathCheckDLL", "SecurityModule")
    hwp.Open(hwp_path)
    sleep(2)
    hwp.Run("SelectAll")
    hwp.Run("Copy")

    win32clipboard.OpenClipboard()
    raw = win32clipboard.GetClipboardData()
    win32clipboard.CloseClipboard()

    hwp.Quit()
    return raw


# def text_copy(raw):
#     make_ppt_text(raw)
#     ppt = Powerpoint()
#     app = ppt.init_app()
#     prs = ppt.open_prs(SAVEFILE_NAME)


def text_copy(raw, src_ppt):
    (
        main_book,
        main_chapter,
        main_verse_start,
        main_verse_end,
        main_title,
    ) = extract_main_verse(raw)
    keys, contentsDict = dict_contents(
        main_book, main_chapter, main_verse_start, main_verse_end
    )
    first_slide, section_count = src_ppt.get_section("성경봉독")
    src_ppt.del_section(first_slide + 1, section_count - 2)
    src_ppt.change_subtitle(
        first_slide,
        f"{bookDict[main_book]} {main_chapter}장 {main_verse_start}-{main_verse_end}절",
    )
    src_ppt.input_verse(first_slide, keys, contentsDict)
    # resultList = copy_contents(raw)
    (
        resultList,
        main_book,
        main_chapter,
        main_verse_start,
        main_verse_end,
        main_title,
    ) = parsing_contents(raw)
    first_slide, section_count = src_ppt.get_section("말씀 선포")
    src_ppt.del_section(first_slide + 1, section_count - 2)
    src_ppt.change_subtitle(first_slide, main_title)
    src_ppt.input_hwp(first_slide, resultList)


# 안쓰는 함수
def text_prs():

    from bible_function import (
        abbreviation,
        parse_paragraph,
        dict_contents,
    )  # 성경 약자 반환 함수 # 정규식 파싱 함수 # 검색 구절을 사전형태로 반환하는 함수

    from pptx import Presentation

    # 파일이름 지정
    template_path = "./verse_template.pptx"
    save_file_name = "new-file-name.pptx"

    # 본문 삽입
    prs = Presentation(template_path)

    # 성경구절 입력
    paragraph = "사도행전 8장 1-8절"

    # 파싱
    main_book, main_chapter, main_verse_start, main_verse_end = parse_paragraph(
        paragraph
    )
    # 리스트 & 딕셔너리 반환
    keys, contentsDict = dict_contents(
        main_book, main_chapter, main_verse_start, main_verse_end
    )

    for key in keys:
        contents_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(contents_slide_layout)
        title = slide.shapes.title
        contents = slide.placeholders[1]

        title.text = key
        contents.text = contentsDict[key]

    prs.save(save_file_name)
    return save_file_name


# 교독문 복사 함수
def responsive_reading_copy(re_no, src_ppt, section_index=3):
    """
    교독문 복사 함수
        section_index : 구역순서 { '교독문' : 3  }
    """
    # DIR = "C:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\교독문_확대"
    # file_name = [file for file in listdir(DIR) if file[:3] == f"{str(re_no).zfill(3)}"][0]
    DIR = "C:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\교독문_16x9"
    file_name = f"교독문{re_no}.pptx"

    print(file_name)
    path = f"{DIR}\\{file_name}"

    print("경로 : ", path)
    ppt = Powerpoint()
    ppt.init_app()
    ppt.open_prs(path)
    ppt.copy_slide_all()
    first_slide, section_count = src_ppt.get_section_number(section_index)
    # src_ppt.prs.Slides.Range(10).Select()
    # src_ppt.prs.Slides.Paste(10)

    print(first_slide)
    sleep(1)
    src_ppt.del_section(first_slide + 1, section_count - 1)
    ### 복사 붙여넣기
    sleep(1)
    print("src_ppt 슬라이드 삽입위치", first_slide)
    src_ppt.app.Windows(2).View.GotoSlide(first_slide)
    src_ppt.app.Windows(2).Activate()  # 이전 창(소스ppt) 활성화
    src_ppt.app.CommandBars.ExecuteMso("PasteSourceFormatting")  # 원본소스유지 붙여넣기
    ###
    # src_ppt.paste_slide(first_slide + 1)
    # ppt.copy_desgin_slide(src_ppt.prs, first_slide + 1)
    ppt.prs.Close()


# 찬양 복사 함수
def hymn_copy(hymn_number, src_ppt, section_index):
    """
    찬양 복사 함수
        section_index : 구역순서 { '찬송 1' : 4,  '찬송2' : 8 }
        hymn_number : 새찬송가 번호
        src_ppt : 소스 ppt
    """
    HYMN_DIR = "C:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\새찬송가16x9"
    hymn_path = f"{HYMN_DIR}\\NHymn{str(hymn_number).zfill(3)}h_Wide.PPT"
    print("찬송가 경로 : ", hymn_path)
    hymn_ppt = Powerpoint()
    hymn_ppt.init_app()
    hymn_ppt.open_prs(hymn_path)
    sleep(0.5)
    hymn_ppt.copy_slide_all()
    print(hymn_ppt.count)
    sleep(0.5)
    first_slide, section_count = src_ppt.get_section_number(section_index)
    src_ppt.del_section(first_slide + 1, section_count - 2)
    win_number = src_ppt.app.Windows.Count  # 현재 ppt 창 번호 따기
    print("win_number", win_number)
    ###
    sleep(1)
    print("src_ppt 슬라이드 삽입위치", first_slide)
    src_ppt.app.Windows(2).View.GotoSlide(first_slide)
    src_ppt.app.Windows(2).Activate()  # 이전 창(소스ppt) 활성화
    src_ppt.app.CommandBars.ExecuteMso("PasteSourceFormatting")  # 원본소스유지 붙여넣기
    ###
    # src_ppt.paste_slide(first_slide + 1)
    sleep(0.5)
    # hymn_ppt.copy_desgin_slide(src_ppt.prs, first_slide + 1)
    src_ppt.change_hymn_number(hymn_number, section_index)
    hymn_ppt.prs.Close()


def bible_prs(raw):
    from bible_function import (
        abbreviation,
        parse_paragraph,
        parsing_contents,
        bookDict,
        dict_contents,
    )

    import re
    from pptx import Presentation
    from pptx.util import Pt

    (
        resultList,
        main_book,
        main_chapter,
        main_verse_start,
        main_verse_end,
        main_title,
    ) = parsing_contents(raw)

    keys, contentsDict = dict_contents(
        main_book, main_chapter, main_verse_start, main_verse_end
    )

    # 파일이름 지정
    template_path = "./verse_template.pptx"
    save_file_name = "new-file-name.pptx"

    # 본문 삽입
    prs = Presentation(template_path)

    for key in keys:
        contents_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(contents_slide_layout)
        title = slide.shapes.title
        contents = slide.placeholders[1]

        title.text = key
        contents.text = contentsDict[key]

    for book, chapter, verse_start, verse_end, contents in resultList:
        if not verse_end:
            key = f"{bookDict[book]} {chapter}장 {verse_start}절"
            # pptx
            contents_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(contents_slide_layout)
            title_placeholder = slide.shapes.title
            contents_placeholder = slide.placeholders[1]

            title_placeholder.text = key
            contents_placeholder.text = contents[0]
        else:
            key = f"{bookDict[book]} {chapter}장 {verse_start}-{verse_end}절"

            contents_slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(contents_slide_layout)

            title_placeholder = slide.shapes.title
            title_placeholder.text = key

            contents_text_frame = slide.placeholders[1].text_frame
            for _ in range(len(contents) - 1):
                contents_text_frame.add_paragraph()
            for i in range(len(contents)):
                tmp1 = contents_text_frame.paragraphs[i].add_run()
                tmp1.text = re.search("\d+", contents[i]).group()
                tmp1.font.size = Pt(36)
                tmp2 = contents_text_frame.paragraphs[i].add_run()
                tmp2.text = re.sub("\d+", "", contents[i])

    prs.save(save_file_name)
    return save_file_name, main_title


def copy_text(raw, src_ppt):
    raw = hwp(hwp_path)
    ######## 정리 필요 ##########
    file_name, main_title = bible_prs(raw)
    print(file_name)

    path = os.getcwd() + os.sep + file_name

    bible_ppt = Powerpoint()
    bible_ppt.init_app()
    bible_ppt.open_prs(path)
    bible_ppt.get_count_slide()
    first_slide, section_count = src_ppt.get_section_number(7)
    bible_ppt.copy_slide_all()
    src_ppt.del_section(first_slide, section_count - 1)
    src_ppt.paste_slide(first_slide + 1)
    bible_ppt.copy_desgin_slide(src_ppt.prs, first_slide)

    src_ppt.prs.Slides(first_slide).Shapes(2).TextFrame.TextRange.Text = main_title
    bible_ppt.prs.Close()
    ######### 정리 필요 #########


def executing():  # 실행
    # 소스 ppt 열기
    path = "c:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\오전예배 (16x9)_2021_base.pptx"
    src_ppt = Powerpoint()
    src_ppt.init_app()
    src_prs = src_ppt.open_prs(path=path)

    re_no = input("교독문 번호를 입력하세요: ")
    responsive_reading_copy(re_no, src_ppt=src_ppt, section_index=3)

    hymn_number = input("첫번째 찬송가 번호를 입력하세요: ")
    hymn_copy(hymn_number, src_ppt=src_ppt, section_index=4)

    raw = hwp(hwp_path)

    hymn_number = input("두번째 찬송가 번호를 입력하세요: ")
    hymn_copy(hymn_number, src_ppt=src_ppt, section_index=8)

    sleep(1)

    src_ppt.change_transition()


if __name__ == "__main__":
    executing()