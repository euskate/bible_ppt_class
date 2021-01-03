import sys
import win32com.client
from time import sleep
from os import listdir
import os

path = (
    "c:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\오전예배 (16x10)_20201227.pptx"
)

hymn1_path = (
    "C:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\새찬송가16x9\\NHymn016h_Wide.PPT"
)


class Powerpoint:
    """
    Python PowerPoint win32 object
    """

    def set_data(self):
        self.path = path
        self.count = 0
        # self

    def init_app(self):
        self.app = win32com.client.Dispatch("PowerPoint.Application")
        self.app.Visible = True
        return self.app

    def quit_app(self):
        self.app.Quit()

    def open_prs(self, path):
        self.prs = self.app.Presentations.Open(path)
        return self.prs

    def get_count_slide(self):
        return len(self.prs.Slides)

    def copy_slide_all(self):
        count = self.get_count_slide()
        self.prs.Slides.Range(range(1, count + 1)).copy()

    def paste_slide(self, paste_slide_number):
        self.prs.Slides.paste(paste_slide_number)

    def copy_desgin_slide(self, dst_prs, desgin_slide_number):
        count = self.get_count_slide()
        for i in range(count):
            dst_prs.Slides[desgin_slide_number + i].Design = self.prs.Slides[
                i + 1
            ].Design

    def get_section_number(self, section_index):
        sec = self.prs.SectionProperties
        self.first_slide = sec.FirstSlide(section_index)
        self.section_count = sec.SlidesCount(section_index)
        return self.first_slide, self.section_count

    def del_section(self, first_slide, section_count):
        self.prs.Slides.Range(
            range(first_slide + 1, first_slide + section_count - 1)
        ).Delete()

    def change_hymn_number(self, hymn_number, section_index):
        first_slide, _ = self.get_section_number(section_index)
        self.prs.Slides(first_slide).Shapes(2).Textframe.TextRange.Runs(
            2
        ).Text = hymn_number
        pass


### 한글 ###
hwp_path = "D:\\Documents\\카카오톡 받은 파일\\01-2. 회개하고 주께 기도하라(ppt용 성경구절).hwp"


def hwp(hwp_path):
    """ 한글함수 """
    import win32clipboard

    hwp = win32com.client.gencache.EnsureDispatch("HWPFrame.HwpObject")
    hwp.RegisterModule("FilePathCheckDLL", "SecurityModule")
    hwp.Open(hwp_path)
    hwp.Run("SelectAll")
    hwp.Run("Copy")

    win32clipboard.OpenClipboard()
    raw = win32clipboard.GetClipboardData()
    win32clipboard.CloseClipboard()

    hwp.Quit()
    return raw


# # 소스 ppt 열기
src_ppt = Powerpoint()
src_app = src_ppt.init_app()
src_prs = src_ppt.open_prs(path)

# def text_prs():

#     from bible_function import (
#         abbreviation,
#         parse_paragraph,
#         dict_contents,
#     )  # 성경 약자 반환 함수 # 정규식 파싱 함수 # 검색 구절을 사전형태로 반환하는 함수

#     from pptx import Presentation

#     # 파일이름 지정
#     template_path = "./verse_template.pptx"
#     save_file_name = "new-file-name.pptx"

#     # 본문 삽입
#     prs = Presentation(template_path)

#     # 성경구절 입력
#     paragraph = "사도행전 8장 1-8절"

#     # 파싱
#     main_book, main_chapter, main_verse_start, main_verse_end = parse_paragraph(
#         paragraph
#     )
#     # 리스트 & 딕셔너리 반환
#     keys, contentsDict = dict_contents(
#         main_book, main_chapter, main_verse_start, main_verse_end
#     )

#     for key in keys:
#         contents_slide_layout = prs.slide_layouts[1]
#         slide = prs.slides.add_slide(contents_slide_layout)
#         title = slide.shapes.title
#         contents = slide.placeholders[1]

#         title.text = key
#         contents.text = contentsDict[key]

#     prs.save(save_file_name)
#     return save_file_name


def responsive_reading_copy(section_index, re_no, src_ppt=src_ppt):
    """
    교독문 복사 함수
        section_index : { '교독문' : 3  }
    """
    DIR = "C:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\교독문_확대"
    file_name = [file for file in listdir(DIR) if file[:3] == f"{str(re_no).zfill(3)}"][
        0
    ]
    print(file_name)
    path = f"{DIR}\\{file_name}"

    print("경로 : ", path)
    ppt = Powerpoint()
    ppt.init_app()
    ppt.open_prs(path)
    ppt.copy_slide_all()
    count = ppt.get_count_slide()
    first_slide, section_count = src_ppt.get_section_number(section_index)
    # src_ppt.prs.Slides.Range(10).Select()
    # src_ppt.prs.Slides.Paste(10)

    print(first_slide)
    sleep(1)
    src_ppt.del_section(first_slide, section_count)
    src_ppt.paste_slide(first_slide + 1)
    ppt.copy_desgin_slide(src_ppt.prs, first_slide)
    ppt.prs.Close()


def hymn_copy(section_index, hymn_number, src_ppt=src_ppt):
    """
    찬양 복사 함수
        section_index : { '찬송 1' : 4,  '찬송2' : 8 }
    """
    HYMN_DIR = "C:\\Users\\Administrator\\Desktop\\WorkSpace\\pyPptx\\새찬송가16x9"
    hymn_path = f"{HYMN_DIR}\\NHymn{str(hymn_number).zfill(3)}h_Wide.PPT"
    print("찬송가 경로 : ", hymn_path)
    hymn_ppt = Powerpoint()
    hymn_ppt.init_app()
    hymn_ppt.open_prs(hymn_path)
    sleep(2)
    hymn_ppt.copy_slide_all()

    first_slide, section_count = src_ppt.get_section_number(section_index)
    print(first_slide)
    src_ppt.del_section(first_slide, section_count)
    src_ppt.paste_slide(first_slide + 1)
    hymn_ppt.copy_desgin_slide(src_ppt.prs, first_slide)
    src_ppt.change_hymn_number(hymn_number, section_index)
    hymn_ppt.prs.Close()


re_no = input("교독문 번호를 입력하세요: ")
responsive_reading_copy(3, re_no)

hymn_number = input("첫번째 찬송가 번호를 입력하세요: ")
hymn_copy(4, hymn_number)


def bible_prs(raw):
    from bible_function import (
        abbreviation,
        parse_paragraph,
        parsing_contents,
        bookDict,
        dict_contents,
    )

    import re
    from pptx import Presentation
    from pptx.util import Pt

    (
        resultList,
        main_book,
        main_chapter,
        main_verse_start,
        main_verse_end,
        main_title,
    ) = parsing_contents(raw)

    keys, contentsDict = dict_contents(
        main_book, main_chapter, main_verse_start, main_verse_end
    )

    # 파일이름 지정
    template_path = "./verse_template.pptx"
    save_file_name = "new-file-name.pptx"

    # 본문 삽입
    prs = Presentation(template_path)

    for key in keys:
        contents_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(contents_slide_layout)
        title = slide.shapes.title
        contents = slide.placeholders[1]

        title.text = key
        contents.text = contentsDict[key]

    for book, chapter, verse_start, verse_end, contents in resultList:
        if not verse_end:
            key = f"{bookDict[book]} {chapter}장 {verse_start}절"
            # pptx
            contents_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(contents_slide_layout)
            title_placeholder = slide.shapes.title
            contents_placeholder = slide.placeholders[1]

            title_placeholder.text = key
            contents_placeholder.text = contents[0]
        else:
            key = f"{bookDict[book]} {chapter}장 {verse_start}-{verse_end}절"

            contents_slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(contents_slide_layout)

            title_placeholder = slide.shapes.title
            title_placeholder.text = key

            contents_text_frame = slide.placeholders[1].text_frame
            for _ in range(len(contents) - 1):
                contents_text_frame.add_paragraph()
            for i in range(len(contents)):
                tmp1 = contents_text_frame.paragraphs[i].add_run()
                tmp1.text = re.search("\d+", contents[i]).group()
                tmp1.font.size = Pt(36)
                tmp2 = contents_text_frame.paragraphs[i].add_run()
                tmp2.text = re.sub("\d+", "", contents[i])

    prs.save(save_file_name)
    return save_file_name, main_title


raw = hwp(hwp_path)
file_name, main_title = bible_prs(raw)
print(file_name)

path = os.getcwd() + os.sep + file_name

bible_ppt = Powerpoint()
bible_ppt.init_app()
bible_ppt.open_prs(path)
count = bible_ppt.get_count_slide()
first_slide, section_count = src_ppt.get_section_number(7)
bible_ppt.copy_slide_all()
src_ppt.del_section(first_slide, section_count)
src_ppt.paste_slide(first_slide + 1)
bible_ppt.copy_desgin_slide(src_ppt.prs, first_slide)


src_prs.Slides(first_slide).Shapes(2).TextFrame.TextRange.Text = main_title
bible_ppt.prs.Close()


hymn_number = input("두번째 찬송가 번호를 입력하세요: ")
hymn_copy(8, hymn_number)


sleep(1)